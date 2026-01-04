"""
年次財務諸表のAI診断レポート生成ビュー
"""
import logging
import json
from typing import Dict, Any
from django.contrib.auth.mixins import LoginRequiredMixin
from django.views import View
from django.http import JsonResponse, HttpResponse
from django.db import transaction
from django.utils import timezone

from ..models import FiscalSummary_Year, Company
from ..mixins import SelectedCompanyMixin, ErrorHandlingMixin
from ..utils.fiscal_ai_diagnosis import (
    collect_fiscal_data_for_diagnosis,
    build_ai_diagnosis_prompt,
)
from ..utils.gemini import get_gemini_response

logger = logging.getLogger(__name__)


class FiscalAIDiagnosisGenerateView(SelectedCompanyMixin, LoginRequiredMixin, ErrorHandlingMixin, View):
    """AI診断レポート生成API"""
    
    def post(self, request, fiscal_summary_year_id):
        try:
            fiscal_summary_year = FiscalSummary_Year.objects.get(
                id=fiscal_summary_year_id,
                company=self.this_company
            )
            
            # 会社情報の不足をチェック
            company = self.this_company
            missing_company_info = []
            company_edit_url = None
            
            if not company.industry_classification:
                missing_company_info.append('業界分類')
            if not company.industry_subclassification:
                missing_company_info.append('業界小分類')
            # company_sizeはCharFieldでデフォルト値's'があるため、Noneになることはない
            # ただし、ユーザーが明示的に設定していない可能性を考慮して、デフォルト値の場合はチェックしない
            # （デフォルト値's'は有効な値として扱う）
            
            # 会社情報が不足している場合
            if missing_company_info:
                from django.urls import reverse
                company_edit_url = reverse('company_update', kwargs={'id': company.id})
                missing_info_text = '、'.join(missing_company_info)
                message = f'診断に必要な情報が不足しています。\n\n以下の項目を設定してください：{missing_info_text}\n\n設定は<a href="{company_edit_url}" target="_blank">会社情報編集画面</a>から行えます。'
                
                return JsonResponse({
                    'success': True,
                    'needs_info': True,
                    'message': message,
                    'missing_info_questions': [f"{missing_info_text}を設定してください。"],
                    'company_edit_url': company_edit_url,
                })
            
            # データを収集
            fiscal_data = collect_fiscal_data_for_diagnosis(
                self.this_company,
                fiscal_summary_year.year
            )
            
            # 不足している情報をチェック
            missing_info = []
            if not fiscal_data['fiscal_data'].get(f'year_{fiscal_summary_year.year}'):
                missing_info.append('対象年度の決算データ')
            if not fiscal_data['fiscal_data'].get(f'year_{fiscal_summary_year.year - 1}'):
                missing_info.append('前期の決算データ')
            if not fiscal_data['fiscal_data'].get(f'year_{fiscal_summary_year.year - 2}'):
                missing_info.append('前々期の決算データ')
            
            # ベンチマークデータが不足している場合（業界分類・企業規模が設定されていても、ベンチマークデータ自体が存在しない場合）
            if not fiscal_data['benchmark_data']:
                if company.industry_classification and company.industry_subclassification and company.company_size:
                    missing_info.append('ローカルベンチマークデータ（該当する業界・規模のベンチマークデータが存在しません）')
                else:
                    # 業界分類・企業規模が未設定の場合は、上記のチェックで既に処理されている
                    pass
            
            # 情報が不足している場合、会話形式で入力
            if missing_info:
                questions = []
                for info in missing_info:
                    if '決算データ' in info:
                        questions.append(f"{info}を入力してください。")
                    elif 'ベンチマーク' in info:
                        questions.append("業界分類と企業規模を設定してください。")
                
                return JsonResponse({
                    'success': True,
                    'needs_info': True,
                    'message': '診断に必要な情報が不足しています。以下の情報を入力してください：',
                    'missing_info_questions': questions,
                })
            
            # プロンプトを構築
            prompt = build_ai_diagnosis_prompt(fiscal_data)
            
            # AI分析を実行
            system_instruction = """あなたは財務分析の専門家です。与えられた財務データとローカルベンチマークデータを基に、詳細で実践的な分析レポートを作成してください。
レポートは3ページ構成で、以下の内容を含めてください：
1. 総合評価と主要指標の分析
2. ローカルベンチマークとの詳細比較
3. 改善提案と今後の展望

分析は具体的で実践的であることを心がけ、数値に基づいた客観的な評価を行ってください。"""
            
            # より軽量なモデルを試す（クォータ制限を回避するため）
            models_to_try = ['gemini-1.5-flash', 'gemini-1.5-pro', 'gemini-2.0-flash-exp']
            report = None
            last_error = None
            
            for model_name in models_to_try:
                try:
                    report = get_gemini_response(
                        prompt=prompt,
                        system_instruction=system_instruction,
                        model=model_name
                    )
                    if report:
                        break
                except ValueError as ve:
                    # クォータ制限エラーの場合は、ユーザーに分かりやすいメッセージを返す
                    if '利用制限' in str(ve) or 'quota' in str(ve).lower():
                        return JsonResponse({
                            'success': False,
                            'error': str(ve)
                        }, status=429)
                    last_error = ve
                    continue
                except Exception as e:
                    logger.warning(f"Error with model {model_name}: {e}")
                    last_error = e
                    continue
            
            if not report:
                error_message = 'AI診断レポートの生成に失敗しました。'
                if last_error:
                    if '利用制限' in str(last_error) or 'quota' in str(last_error).lower():
                        error_message = str(last_error)
                    else:
                        error_message = f'AI診断レポートの生成に失敗しました: {str(last_error)}'
                
                return JsonResponse({
                    'success': False,
                    'error': error_message
                }, status=500)
            
            return JsonResponse({
                'success': True,
                'needs_info': False,
                'report': report,
            })
            
        except FiscalSummary_Year.DoesNotExist:
            return JsonResponse({
                'success': False,
                'error': '決算データが見つかりません。'
            }, status=404)
        except ValueError as ve:
            # クォータ制限エラーなど、ValueErrorとして処理されるエラー
            if '利用制限' in str(ve) or 'quota' in str(ve).lower():
                return JsonResponse({
                    'success': False,
                    'error': str(ve)
                }, status=429)
            return JsonResponse({
                'success': False,
                'error': str(ve)
            }, status=400)
        except Exception as e:
            logger.error(f"Error generating AI diagnosis report: {e}", exc_info=True)
            return JsonResponse({
                'success': False,
                'error': f'エラーが発生しました: {str(e)}'
            }, status=500)


class FiscalAIDiagnosisChatView(SelectedCompanyMixin, LoginRequiredMixin, ErrorHandlingMixin, View):
    """AI診断の会話形式情報入力API"""
    
    def post(self, request, fiscal_summary_year_id):
        try:
            import json
            data = json.loads(request.body)
            user_message = data.get('message', '')
            
            fiscal_summary_year = FiscalSummary_Year.objects.get(
                id=fiscal_summary_year_id,
                company=self.this_company
            )
            
            # 会社情報の不足をチェック（初回リクエスト時のみ）
            company = self.this_company
            if not user_message:  # 初回リクエストの場合
                missing_company_info = []
                if not company.industry_classification:
                    missing_company_info.append('業界分類')
                if not company.industry_subclassification:
                    missing_company_info.append('業界小分類')
                
                if missing_company_info:
                    from django.urls import reverse
                    company_edit_url = reverse('company_update', kwargs={'id': company.id})
                    missing_info_text = '、'.join(missing_company_info)
                    message = f'診断に必要な情報が不足しています。\n\n以下の項目を設定してください：{missing_info_text}\n\n設定は<a href="{company_edit_url}" target="_blank">会社情報編集画面</a>から行えます。'
                    
                    return JsonResponse({
                        'success': True,
                        'needs_info': True,
                        'response': message,
                        'company_edit_url': company_edit_url,
                    })
            
            # 会話履歴を取得（セッションから）
            chat_history = request.session.get(f'diagnosis_chat_{fiscal_summary_year_id}', [])
            chat_history.append({'role': 'user', 'content': user_message})
            
            # AI応答を生成
            system_instruction = """あなたは財務分析の専門家です。ユーザーから不足している財務情報を聞き出し、適切な質問をして情報を収集してください。
情報が十分に集まったら、分析レポートを作成してください。"""
            
            # 会話履歴をプロンプトに含める
            conversation = "\n".join([f"{msg['role']}: {msg['content']}" for msg in chat_history])
            prompt = f"以下の会話を続けてください：\n\n{conversation}\n\nassistant:"
            
            ai_response = get_gemini_response(
                prompt=prompt,
                system_instruction=system_instruction,
                model='gemini-2.0-flash-exp'
            )
            
            chat_history.append({'role': 'assistant', 'content': ai_response})
            request.session[f'diagnosis_chat_{fiscal_summary_year_id}'] = chat_history
            
            # レポートが準備できたかどうかを判定（簡易版：会話が一定回数以上の場合）
            report_ready = len(chat_history) >= 6  # 3往復以上
            
            response_data = {
                'success': True,
                'response': ai_response,
                'report_ready': report_ready,
            }
            
            if report_ready:
                # レポートを生成
                fiscal_data = collect_fiscal_data_for_diagnosis(
                    self.this_company,
                    fiscal_summary_year.year
                )
                prompt = build_ai_diagnosis_prompt(fiscal_data, missing_info=None)
                report = get_gemini_response(
                    prompt=prompt,
                    system_instruction=system_instruction,
                    model='gemini-2.0-flash-exp'
                )
                response_data['report'] = report
            
            return JsonResponse(response_data)
            
        except FiscalSummary_Year.DoesNotExist:
            return JsonResponse({
                'success': False,
                'error': '決算データが見つかりません。'
            }, status=404)
        except Exception as e:
            logger.error(f"Error in AI diagnosis chat: {e}", exc_info=True)
            return JsonResponse({
                'success': False,
                'error': f'エラーが発生しました: {str(e)}'
            }, status=500)


class FiscalAIDiagnosisDownloadView(SelectedCompanyMixin, LoginRequiredMixin, ErrorHandlingMixin, View):
    """AI診断レポートのダウンロードビュー"""
    
    def get(self, request, fiscal_summary_year_id, format_type):
        try:
            fiscal_summary_year = FiscalSummary_Year.objects.get(
                id=fiscal_summary_year_id,
                company=self.this_company
            )
            
            # プランによる制限チェック
            can_download_advanced = False
            try:
                if hasattr(request.user, 'userfirm') and request.user.userfirm.exists():
                    user_firm = request.user.userfirm.filter(is_selected=True, active=True).first()
                    if user_firm:
                        subscription = user_firm.firm.subscription
                        plan_type = subscription.plan.plan_type
                        can_download_advanced = plan_type in ['professional', 'enterprise']
            except Exception:
                pass
            
            # フォーマットチェック
            if format_type not in ['pdf', 'pptx', 'excel']:
                return JsonResponse({'error': '無効なフォーマットです。'}, status=400)
            
            if format_type in ['pptx', 'excel'] and not can_download_advanced:
                return JsonResponse({
                    'error': 'このフォーマットはProfessionalプラン以上で利用できます。'
                }, status=403)
            
            # レポートを生成（セッションから取得、または再生成）
            report_text = request.session.get(f'diagnosis_report_{fiscal_summary_year_id}')
            if not report_text:
                # レポートがセッションにない場合、再生成
                fiscal_data = collect_fiscal_data_for_diagnosis(
                    self.this_company,
                    fiscal_summary_year.year
                )
                prompt = build_ai_diagnosis_prompt(fiscal_data)
                system_instruction = """あなたは財務分析の専門家です。詳細で実践的な分析レポートを作成してください。"""
                report_text = get_gemini_response(
                    prompt=prompt,
                    system_instruction=system_instruction,
                    model='gemini-2.0-flash-exp'
                )
            
            # フォーマットに応じてレポートを生成
            if format_type == 'pdf':
                return self._generate_pdf(report_text, fiscal_summary_year)
            elif format_type == 'pptx':
                return self._generate_pptx(report_text, fiscal_summary_year)
            elif format_type == 'excel':
                return self._generate_excel(report_text, fiscal_summary_year)
            
        except FiscalSummary_Year.DoesNotExist:
            return JsonResponse({
                'error': '決算データが見つかりません。'
            }, status=404)
        except Exception as e:
            logger.error(f"Error downloading AI diagnosis report: {e}", exc_info=True)
            return JsonResponse({
                'error': f'エラーが発生しました: {str(e)}'
            }, status=500)
    
    def _generate_pdf(self, report_text: str, fiscal_summary_year: FiscalSummary_Year) -> HttpResponse:
        """PDFレポートを生成"""
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.enums import TA_LEFT, TA_CENTER
            import io
            
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
            
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=16,
                textColor=colors.HexColor('#366092'),
                spaceAfter=30,
                alignment=TA_CENTER
            )
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontSize=12,
                textColor=colors.HexColor('#366092'),
                spaceAfter=12,
                spaceBefore=12
            )
            
            story = []
            
            # タイトル
            title = Paragraph(f"{fiscal_summary_year.company.name} - AI診断レポート（{fiscal_summary_year.year}年度）", title_style)
            story.append(title)
            story.append(Spacer(1, 0.2*inch))
            
            # レポート内容をパースして追加
            lines = report_text.split('\n')
            current_page = 1
            for line in lines:
                line = line.strip()
                if not line:
                    story.append(Spacer(1, 0.1*inch))
                elif line.startswith('###'):
                    # 小見出し
                    text = line.replace('###', '').strip()
                    story.append(Paragraph(text, heading_style))
                elif line.startswith('##'):
                    # 中見出し
                    if current_page > 1:
                        story.append(PageBreak())
                    text = line.replace('##', '').strip()
                    story.append(Paragraph(text, heading_style))
                    current_page += 1
                elif line.startswith('#'):
                    # 大見出し
                    if current_page > 1:
                        story.append(PageBreak())
                    text = line.replace('#', '').strip()
                    story.append(Paragraph(text, title_style))
                    current_page += 1
                else:
                    # 本文
                    story.append(Paragraph(line, styles['Normal']))
            
            doc.build(story)
            
            response = HttpResponse(buffer.getvalue(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="ai_diagnosis_{fiscal_summary_year.company.name}_{fiscal_summary_year.year}.pdf"'
            return response
            
        except ImportError:
            return JsonResponse({'error': 'PDF生成ライブラリ（reportlab）がインストールされていません。'}, status=500)
    
    def _generate_pptx(self, report_text: str, fiscal_summary_year: FiscalSummary_Year) -> HttpResponse:
        """PPTXレポートを生成"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import io
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # レポートを3ページに分割
            lines = report_text.split('\n')
            pages = []
            current_page = []
            
            for line in lines:
                if line.startswith('##') and current_page:
                    pages.append('\n'.join(current_page))
                    current_page = [line]
                else:
                    current_page.append(line)
            if current_page:
                pages.append('\n'.join(current_page))
            
            # 各ページをスライドとして追加
            for i, page_content in enumerate(pages[:3]):  # 最大3ページ
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白レイアウト
                
                # タイトル
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = f"AI診断レポート - ページ {i+1}"
                title_frame.paragraphs[0].font.size = Pt(24)
                title_frame.paragraphs[0].font.bold = True
                
                # 本文
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                # 内容を追加
                for line in page_content.split('\n'):
                    if line.strip():
                        p = content_frame.add_paragraph()
                        p.text = line.strip()
                        p.font.size = Pt(12)
            
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            response = HttpResponse(buffer.getvalue(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = f'attachment; filename="ai_diagnosis_{fiscal_summary_year.company.name}_{fiscal_summary_year.year}.pptx"'
            return response
            
        except ImportError:
            return JsonResponse({'error': 'PPTX生成ライブラリ（python-pptx）がインストールされていません。'}, status=500)
    
    def _generate_excel(self, report_text: str, fiscal_summary_year: FiscalSummary_Year) -> HttpResponse:
        """Excelレポートを生成"""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment
            import io
            
            wb = Workbook()
            ws = wb.active
            ws.title = "AI診断レポート"
            
            # タイトル
            ws['A1'] = f"{fiscal_summary_year.company.name} - AI診断レポート（{fiscal_summary_year.year}年度）"
            ws['A1'].font = Font(size=16, bold=True)
            ws.merge_cells('A1:D1')
            
            # レポート内容を追加
            row = 3
            lines = report_text.split('\n')
            for line in lines:
                if line.strip():
                    ws[f'A{row}'] = line.strip()
                    if line.startswith('#'):
                        ws[f'A{row}'].font = Font(bold=True, size=12)
                    row += 1
            
            # 列幅を調整
            ws.column_dimensions['A'].width = 100
            
            buffer = io.BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            
            response = HttpResponse(buffer.getvalue(), content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
            response['Content-Disposition'] = f'attachment; filename="ai_diagnosis_{fiscal_summary_year.company.name}_{fiscal_summary_year.year}.xlsx"'
            return response
            
        except ImportError:
            return JsonResponse({'error': 'Excel生成ライブラリ（openpyxl）がインストールされていません。'}, status=500)

